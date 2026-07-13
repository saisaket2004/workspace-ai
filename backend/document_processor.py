"""
document_processor.py — Text extraction from document files.

WHY:  Centralises document parsing so drive_service.py doesn't need to
      know how to read PDFs, DOCX files, etc.

WHERE: Called by drive_service.py.

HOW:  `from backend.document_processor import extract_pdf_text, extract_docx_text`
"""

from __future__ import annotations

import logging

logger = logging.getLogger(__name__)


def extract_pdf_text(pdf_path: str) -> str:
    """Extract text from a PDF file using PyMuPDF (fitz)."""
    import fitz

    document = fitz.open(pdf_path)
    text = ""
    for page in document:
        text += page.get_text()
    document.close()
    return text


def extract_docx_text(docx_path: str) -> str:
    """
    Extract text from a .docx file using python-docx.

    Raises ImportError if python-docx is not installed.
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "python-docx is required to read .docx files. "
            "Install it with: pip install python-docx"
        )

    doc = Document(docx_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)

def extract_pptx_text(pptx_path: str) -> str:
    """
    Extract text from a .pptx file using python-pptx.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required to read .pptx files. "
            "Install it with: pip install python-pptx"
        )
        
    prs = Presentation(pptx_path)
    text_content = []
    
    for i, slide in enumerate(prs.slides):
        text_content.append(f"--- Slide {i+1} ---")
        
        # Extract shapes text
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
                
        # Extract notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                text_content.append(f"[Speaker Notes]: {notes}")
                
    return "\n".join(text_content)